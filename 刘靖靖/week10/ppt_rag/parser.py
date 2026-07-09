"""PPT 解析：提取每页文本 / 表格 / 备注 / 图片(可选 OCR)。

边界处理：
- 文件不存在 / 损坏 / 加密
- 空页、shape 异常、组合形状递归
- 图片过大或过小
- OCR 失败 / 模型未加载
- 单页文本过长截断
"""
from __future__ import annotations

import os
import tempfile
from pathlib import Path
from typing import Dict, List, Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.exceptions import PackageNotFoundError, InvalidPackageError

from config import OCRConfig
from utils import get_logger

logger = get_logger()

# 单元格内容最长保留长度，避免极端表格撑爆上下文
MAX_CELL_TEXT = 200
# 单页文本最长保留长度（解析阶段先粗截，后续切分再细切）
MAX_PAGE_TEXT = 20000


class ParseError(RuntimeError):
    """PPT 解析失败。"""


def _clean(text: str) -> str:
    """去除空白字符与控制字符，保留可读文本。"""
    if not text:
        return ""
    cleaned = "".join(ch for ch in text if ch == "\t" or ch == "\n" or ord(ch) >= 32)
    lines = [ln.strip() for ln in cleaned.splitlines()]
    return "\n".join(ln for ln in lines if ln)


def _extract_shape(shape, image_collector: Optional[List[bytes]] = None) -> str:
    """递归提取 shape 中的文本与表格，支持组合形状。

    若 image_collector 不为 None，会把图片二进制追加进去（用于后续 OCR）。
    """
    parts: List[str] = []

    # 组合形状：递归子形状
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP and hasattr(shape, "shapes"):
        for sub in shape.shapes:
            parts.append(_extract_shape(sub, image_collector))
        return "\n".join(p for p in parts if p)

    if getattr(shape, "has_text_frame", False):
        parts.append(_clean(shape.text_frame.text))

    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [_clean(c.text)[:MAX_CELL_TEXT] for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))

    # 提取图片二进制（不解析，由 parser 决定是否 OCR）
    if image_collector is not None and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            img = shape.image
            blob = img.blob
            if blob:
                image_collector.append(blob)
        except Exception as e:
            logger.warning(f"图片提取失败，已跳过: {e}")

    return "\n".join(p for p in parts if p)


def _save_image_blob(blob: bytes, tmp_dir: Path, idx: int) -> Optional[Path]:
    """把图片二进制写入临时文件，返回路径。"""
    if not blob:
        return None
    # 通过魔数判断后缀（兜底用 png）
    ext = "png"
    if blob[:3] == b"\xff\xd8\xff":
        ext = "jpg"
    elif blob[:8] == b"\x89PNG\r\n\x1a\n":
        ext = "png"
    elif blob[:4] == b"RIFF" and blob[8:12] == b"WEBP":
        ext = "webp"
    elif blob[:6] in (b"GIF87a", b"GIF89a"):
        ext = "gif"
    elif blob[:2] == b"BM":
        ext = "bmp"

    try:
        tmp_dir.mkdir(parents=True, exist_ok=True)
        path = tmp_dir / f"img_{idx:04d}.{ext}"
        with open(path, "wb") as f:
            f.write(blob)
        return path
    except OSError as e:
        logger.warning(f"临时图片写入失败: {e}")
        return None


def parse_pptx(path, ocr_cfg: Optional[OCRConfig] = None) -> List[Dict]:
    """解析 PPT 文件。

    返回 [{"page": int, "content": str, "source": str}]，跳过空页。
    若 ocr_cfg.enabled=True，会对每页图片做 OCR 并合并到 content。
    抛出 ParseError 表示文件级失败。
    """
    p = Path(path)
    if not p.exists():
        raise ParseError(f"PPT 文件不存在: {p}")
    if p.stat().st_size == 0:
        raise ParseError(f"PPT 文件大小为 0: {p}")

    try:
        prs = Presentation(str(p))
    except (PackageNotFoundError, InvalidPackageError) as e:
        raise ParseError(f"无法打开 PPT（可能损坏或加密）: {p} - {e}")
    except Exception as e:
        raise ParseError(f"解析 PPT 失败: {p} - {e}")

    if not prs.slides:
        logger.warning(f"PPT 无任何 slide: {p}")
        return []

    # 初始化 OCR
    ocr = None
    if ocr_cfg and ocr_cfg.enabled:
        from ocr import get_ocr
        ocr = get_ocr(ocr_cfg)
        if ocr is None:
            logger.warning("OCR 已启用但实例创建失败，将跳过图片文字识别")

    # 临时目录
    tmp_dir = None
    if ocr is not None:
        tmp_dir = Path(ocr_cfg.tmp_dir) / Path(p).stem
        try:
            tmp_dir.mkdir(parents=True, exist_ok=True)
        except OSError as e:
            logger.warning(f"OCR 临时目录创建失败: {e}")
            tmp_dir = None

    pages: List[Dict] = []
    for idx, slide in enumerate(prs.slides, start=1):
        try:
            text_parts: List[str] = []
            image_blobs: List[bytes] = []

            # 提取文本/表格 + 收集图片
            for shape in slide.shapes:
                try:
                    txt = _extract_shape(shape, image_collector=image_blobs)
                    if txt:
                        text_parts.append(txt)
                except Exception as e:
                    logger.warning(f"第 {idx} 页某 shape 提取失败，已跳过: {e}")

            # 备注
            if slide.has_notes_slide:
                try:
                    note = _clean(slide.notes_slide.notes_text_frame.text)
                    if note:
                        text_parts.append(f"[备注] {note}")
                except Exception as e:
                    logger.warning(f"第 {idx} 页备注提取失败: {e}")

            # OCR 图片
            if ocr is not None and image_blobs and tmp_dir is not None:
                ocr_texts = _ocr_slide_images(ocr, image_blobs, tmp_dir, idx, ocr_cfg)
                if ocr_texts:
                    text_parts.append(f"[图片文字]\n{ocr_texts}")

            content = "\n".join(text_parts).strip()
            if len(content) > MAX_PAGE_TEXT:
                logger.warning(f"第 {idx} 页文本过长 ({len(content)})，已截断至 {MAX_PAGE_TEXT}")
                content = content[:MAX_PAGE_TEXT]

            if not content:
                logger.info(f"第 {idx} 页无文本，跳过")
                continue
            pages.append({"page": idx, "content": content, "source": str(p)})
        except Exception as e:
            logger.warning(f"第 {idx} 页整体解析失败，已跳过: {e}")

    logger.info(f"解析完成: {p.name} 共 {len(prs.slides)} 页，有效页 {len(pages)} 页")
    if not pages:
        logger.warning("PPT 解析后无任何有效文本")
    return pages


def _ocr_slide_images(ocr, image_blobs: List[bytes], tmp_dir: Path,
                      slide_idx: int, ocr_cfg: OCRConfig) -> str:
    """对单页的所有图片做 OCR，合并文本。"""
    image_paths = []
    for i, blob in enumerate(image_blobs):
        # 过小图片跳过（图标/装饰）
        if len(blob) < ocr_cfg.min_image_bytes:
            continue
        p = _save_image_blob(blob, tmp_dir / f"slide{slide_idx:03d}", i)
        if p is not None:
            image_paths.append(p)

    if not image_paths:
        return ""

    logger.info(f"第 {slide_idx} 页: 待 OCR 图片 {len(image_paths)} 张")

    # 多张走批量接口，单张走单图接口
    try:
        if len(image_paths) == 1:
            text = ocr.ocr_image(image_paths[0])
        else:
            text = ocr.ocr_images([str(p) for p in image_paths])
    except Exception as e:
        logger.warning(f"第 {slide_idx} 页 OCR 失败: {e}")
        return ""

    if not text:
        logger.info(f"第 {slide_idx} 页 OCR 无识别结果")
        return ""

    # 清理临时图
    for p in image_paths:
        try:
            os.remove(p)
        except OSError:
            pass

    return text.strip()
